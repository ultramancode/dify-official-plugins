"""Abstract interface for pptx document loader implementations."""

import logging
import mimetypes
import re
import uuid
from io import BytesIO
from urllib.parse import urlparse

import requests
from dify_plugin import Tool
from pptx import Presentation

from tools.extractor_base import BaseExtractor
from tools.document import Document, ExtractorResult

logger = logging.getLogger(__name__)


class PPTXExtractor(BaseExtractor):
    """Load pptx files.

    Args:
        tool: Tool instance
        file_bytes: file bytes
        file_name: file name
    """

    def __init__(self, tool: Tool, file_bytes: bytes, file_name: str):
        """Initialize with file path."""
        self._file_bytes = file_bytes
        self._file_name = file_name
        self._tool = tool

    def extract(self) -> ExtractorResult:
        """Load given path as single page."""
        content, img_list = self.parse_pptx(self._file_bytes)
        return ExtractorResult(
            md_content=content,
            documents=[
                Document(page_content=content, metadata={"source": self._file_name})
            ],
            img_list=img_list,
            origin_result=None
        )

    @staticmethod
    def _is_valid_url(url: str) -> bool:
        """Check if the url is valid."""
        parsed = urlparse(url)
        return bool(parsed.netloc) and bool(parsed.scheme)

    def _extract_images_from_pptx(self, prs):
        image_map = {}
        img_list = []
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image_ext = shape.image.ext
                    if image_ext is None:
                        continue
                    file_uuid = str(uuid.uuid4())
                    file_name = file_uuid + "." + image_ext
                    mime_type, _ = mimetypes.guess_type(file_name)

                    file_res = self._tool.session.file.upload(
                        file_name, shape.image.blob, mime_type
                    )

                    image_map[(slide_idx, shape.shape_id)] = f"![image]({file_res.preview_url})"
                    img_list.append(file_res)

        return image_map, img_list

    def _table_to_markdown(self, table):
        markdown = ""
        rows = table.rows
        cols = table.columns
        total_cols = len(cols)

        # Header
        header_cells = [cell.text.strip() for cell in rows[0].cells]
        markdown += "| " + " | ".join(header_cells) + " |\n"
        markdown += "| " + " | ".join(["---"] * total_cols) + " |\n"

        # Rows
        for row in list(rows)[1:]:
            row_cells = [cell.text.strip() for cell in row.cells]
            markdown += "| " + " | ".join(row_cells) + " |\n"

        return markdown

    def parse_pptx(self, file_bytes):
        prs = Presentation(BytesIO(file_bytes))

        content = []

        image_map, img_list = self._extract_images_from_pptx(prs)

        url_pattern = re.compile(r"http://[^\s+]+//|https://[^\s+]+")
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = []
            for shape in slide.shapes:
                # Extract text (including hyperlinks)
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = ""
                        for run in paragraph.runs:
                            run_text = run.text or ""
                            # Check for hyperlink
                            if run.hyperlink and run.hyperlink.address:
                                if url_pattern.match(run.hyperlink.address):
                                    para_text += f"[{run_text}]({run.hyperlink.address})"
                                else:
                                    para_text += run_text
                            else:
                                para_text += run_text
                        if para_text.strip():
                            slide_content.append(para_text.strip())
                # Extract images
                if hasattr(shape, "image"):
                    image_md = image_map.get((slide_idx, shape.shape_id))
                    if image_md:
                        slide_content.append(image_md)
                # Extract tables
                if shape.has_table:
                    table_md = self._table_to_markdown(shape.table)
                    slide_content.append(table_md)
            if slide_content:
                content.append(f"# Slide {slide_idx + 1}\n" + "\n".join(slide_content))
        return "\n\n".join(content), img_list
